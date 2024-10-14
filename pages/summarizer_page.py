import streamlit as st
import os
import tempfile
import sys
import asyncio
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Add the parent directory to sys.path to allow importing from modules
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from modules.summarizer import PDFExtract

def create_ppt(summaries, output_dir, filename):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = os.path.splitext(filename)[0]
    subtitle.text = "Document Summary"

    # Add summaries to slides
    for summary in summaries:
        slide_layout = prs.slide_layouts[1]  # Use layout with title and content
        slide = prs.slides.add_slide(slide_layout)
        
        lines = summary.split('\n')
        
        # Handle the topic line
        if lines[0].startswith("Main Topic:"):
            topic = lines[0].replace("Main Topic:", "Topic:").strip()
            title = slide.shapes.title
            title.text = topic
            lines = lines[1:]  # Remove the topic line from further processing
        
        # Add content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()  # Clear existing text
        
        for line in lines:
            p = tf.add_paragraph()
            p.text = line.strip()
            p.level = 0 if line.startswith('•') else 1
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0, 0, 0)

    # Save the presentation
    ppt_path = os.path.join(output_dir, f"{os.path.splitext(filename)[0]}_summary.pptx")
    prs.save(ppt_path)
    return ppt_path

async def process_pdf(pdf_extract, tmp_file_path, output_dir, unzip_dir, chunked_dir):
    with st.spinner("Parsing PDF..."):
        pdf_extract.parse_pdf(tmp_file_path, os.path.join(output_dir, "extracted.zip"), unzip_dir, chunked_dir)
    st.success("PDF parsed successfully!")

    files = pdf_extract.get_files_from_dir(chunked_dir)
    list_of_all_docs = []
    for file in files:
        if file.endswith(".txt"):
            document = pdf_extract.load_docs(file)
            list_of_all_docs.extend(document)

    with st.spinner("Generating summaries..."):
        summaries = await pdf_extract.process_documents(list_of_all_docs, output_dir)
    st.success("Summaries generated successfully!")

    return summaries

async def show_summarizer_page():
    st.header("Summarizer & Slide Deck Creator")

    uploaded_file = st.file_uploader("Choose a PDF file", type="pdf")

    if uploaded_file is not None:
        st.success("File uploaded successfully!")

        if st.button("Process and Summarize"):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                tmp_file_path = tmp_file.name

            output_dir = tempfile.mkdtemp()
            unzip_dir = os.path.join(output_dir, "unzipped")
            chunked_dir = os.path.join(output_dir, "chunks")
            os.makedirs(unzip_dir)
            os.makedirs(chunked_dir)

            pdf_extract = PDFExtract(os.getenv("PDF_SERVICES_CLIENT_ID"), os.getenv("PDF_SERVICES_CLIENT_SECRET"))

            summaries = await process_pdf(pdf_extract, tmp_file_path, output_dir, unzip_dir, chunked_dir)

            with st.spinner("Creating PowerPoint summary..."):
                ppt_path = create_ppt(summaries, output_dir, uploaded_file.name)
            
            st.success("PowerPoint summary created successfully!")
            st.download_button(
                label="Download PowerPoint Summary",
                data=open(ppt_path, "rb").read(),
                file_name=f"{os.path.splitext(uploaded_file.name)[0]}_summary.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

            # Clean up the temporary file
            os.unlink(tmp_file_path)

    else:
        st.info("Please upload a PDF file to begin.")

def main():
    asyncio.run(show_summarizer_page())

if __name__ == "__main__":
    main()